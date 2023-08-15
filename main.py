import collections
import collections.abc
import pptx
from pptx import Presentation
import os
from PIL import Image
from pptx.util import Inches


class CreatePresentation:

    def __init__(self):
        self.prs = Presentation()

    def _get_slide_layout(self, layout_name):
        """
        The slide layouts in a standard PowerPoint theme always occur in the same sequence.
        This allows content from one deck to be pasted into another and be connected with the right new slide layout:
            * Title (presentation title slide - with subtitle)
            * Title and Content
            * Section Header (sometimes called Segue)
            * Two Content (side by side bullet textboxes)
            * Comparison (same but additional title for each side by side content box)
            * Title Only
            * Blank
            * Content with Caption
            * Picture with Caption
        In python-pptx, these are prs.slide_layouts[0] through prs.slide_layouts[8].
        However, there’s no rule they have to appear in this order, it’s just a convention followed by the themes
        provided with PowerPoint. If the deck you’re using as your template has different slide layouts or has them in
        a different order, you’ll have to work out the slide layout indices for yourself. It’s pretty easy. Just open
        it up in Slide Master view in PowerPoint and count down from the top, starting at zero.
        """
        match layout_name:
            case "Title":
                return self.prs.slide_layouts[0]
            case "Title and Content":
                return self.prs.slide_layouts[1]
            case "Section Header":
                return self.prs.slide_layouts[2]
            case "Two Content":
                return self.prs.slide_layouts[3]
            case "Comparison":
                return self.prs.slide_layouts[4]
            case "Title Only":
                return self.prs.slide_layouts[5]
            case "Blank":
                return self.prs.slide_layouts[6]
            case "Content with Caption":
                return self.prs.slide_layouts[7]
            case "Picture with Caption":
                return self.prs.slide_layouts[8]
            case _:
                raise f"Layout {layout_name} not found!"

    def _add_slide(self, layout_name):
        layout = self._get_slide_layout(layout_name)
        slide = self.prs.slides.add_slide(layout)
        return slide

    def _check_slide_shapes(self, slide, layout_name):
        for shape in slide.shapes:
            print(f"Shape type: {shape.shape_type}")
            if shape.is_placeholder:
                phf = shape.placeholder_format
                print(f"\tName: {shape.name}\n"
                      f"\tIndex: {phf.idx}\n"
                      f"\tType: {phf.type}")

    def add_slide__Picture_w_Caption(self, picture, title='', text='', print_slide_shapes=False):
        layout_name = "Picture with Caption"
        slide = self._add_slide(layout_name=layout_name)

        if print_slide_shapes:
            self._check_slide_shapes(slide, layout_name=layout_name)

        # fill placeholders:
        if title:
            slide.placeholders[0].text = title
        else:
            ph = slide.placeholders[0]
            elem = ph.element
            elem.getparent().remove(elem)

        # todo - adjust/resize placeholder to be bigger (keep img ratio)
        pic_width, pic_height = Image.open(picture).size
        slide.placeholders[1].height = pic_height
        slide.placeholders[1].width = pic_width
        slide.placeholders[1].insert_picture(picture)

        if text:
            slide.placeholders[2].text = text
        else:
            ph = slide.placeholders[2]
            elem = ph.element
            elem.getparent().remove(elem)

    def add_slide__Title(self, title, subtitle='', print_slide_shapes=False):
        layout_name = "Title"
        slide = self._add_slide(layout_name=layout_name)

        if print_slide_shapes:
            self._check_slide_shapes(slide, layout_name=layout_name)

        # fill placeholders:
        slide.placeholders[0].text = title

        if subtitle:
            slide.placeholders[1].text = subtitle
        else:
            ph = slide.placeholders[1]
            elem = ph.element
            elem.getparent().remove(elem)

    def add_slide__Blank(self, add_picture=""):
        layout_name = "Blank"
        slide = self._add_slide(layout_name=layout_name)

        if add_picture:
            slide.shapes.add_picture(add_picture,
                                     left=pptx.util.Inches(0),
                                     top=pptx.util.Inches(0),
                                     width=self.prs.slide_width)

    def save_pptx(self, output_file, start_file=False):
        self.prs.save(output_file)
        if start_file:
            os.startfile(output_file)


if __name__ == "__main__":
    folder = "\\\\ngdp-26\\Integration tests\\Zohar Gazi\\Slowness Investigation\\10_08_23 baseline test"
    # sub_folders = os.listdir(folder)
    # sub_folders = ['NZK2021W41A001', 'NZK2021W41A017', 'YUC2022W32A219', 'YUC2022W32A272', 'YUC2022W32A330', 'YUC2022W32A331']
    sub_folders = ['NZK2021W41A001', "NZK2021W41A017", "YUC2022W32A272", "YUC2022W32A330"]

    for sf in sub_folders:
        for i in range(1, 21):
            print(f"Create Presentation for {sf} - Run {i}")
            presentation = CreatePresentation()
            presentation.add_slide__Title(title=sf, subtitle=f'Run {i}')

            # add uxStats figure slide
            uxStats = os.path.join(folder, sf, str(i), "Last Session Log Analysis", "figures", "uxStats.png")
            presentation.add_slide__Blank(add_picture=uxStats)

            # add DataCollector figures
            DataCollector_figures = os.path.join(folder, sf, str(i), "Data Collector Analysis", "figures")
            for im in os.listdir(DataCollector_figures):
                if im.endswith("png"):
                    fig_path = os.path.join(DataCollector_figures, im)
                    presentation.add_slide__Blank(add_picture=fig_path)

            print("Saving pptx")
            presentation.save_pptx(f"{sf} - {i}.pptx", start_file=False)
            print("Done!")
